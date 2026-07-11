"""Report rendering: interactive HTML detailed analysis + PPT overview.

Both outputs visually separate fact / estimate / assumption everywhere a number
appears (spec: RESEARCH_AGENT_SPEC.md §7). The HTML is fully self-contained —
inline CSS/JS, no CDN. PPT needs python-pptx and is skipped with a warning if missing.
"""
from __future__ import annotations

import html
from pathlib import Path
from typing import List

from src.state.casefile import CaseFile, ToolOption

KIND_COLORS = {"fact": "#1a7f37", "estimate": "#9a6700", "assumption": "#cf222e"}


def _esc(s: object) -> str:
    return html.escape(str(s if s is not None else ""))


def _badge(kind: str) -> str:
    color = KIND_COLORS.get(kind, "#57606a")
    return f'<span class="badge" style="background:{color}">{_esc(kind)}</span>'


def _workflow_rows(steps) -> str:
    rows = []
    for s in steps:
        label = f'<span class="label">{_esc(s.label)}</span>' if s.label else ""
        rows.append(
            f"<tr><td>{_esc(s.id)}</td><td>{_esc(s.name)}</td><td>{_esc(s.actor)}</td>"
            f"<td>{_esc(s.system)}</td><td>{_esc(s.time_estimate)}</td>"
            f"<td>{_esc('; '.join(s.pain_points))}</td><td>{label} {_esc(s.rationale)}</td></tr>"
        )
    return "".join(rows)


def _option_card(o: ToolOption) -> str:
    sim = o.similarity
    exists = ('<p class="exists">Already exists: '
              f'<a href="{_esc(sim.existing_solution_url or o.url)}">'
              f'{_esc(sim.existing_solution_url or o.url)}</a> — gaps: '
              f'{_esc(", ".join(sim.missing) or "none")}</p>') if sim.existing_solution else ""
    matched = "".join(f"<li>{_esc(m)}</li>" for m in sim.matched) or "<li>—</li>"
    missing = "".join(f"<li>{_esc(m)}</li>" for m in sim.missing) or "<li>—</li>"
    assumptions = "".join(f"<li>{_esc(a)}</li>" for a in o.costs.assumptions) or "<li>—</li>"
    scores = ", ".join(f"{_esc(k)}: {v:g}" for k, v in o.scores.items())
    return f"""
<div class="card">
  <h4>{_esc(o.name)} <small>{_esc(o.vendor)}</small>
      <span class="sim">similarity {sim.index}/100</span></h4>
  <p><a href="{_esc(o.url)}">{_esc(o.url)}</a></p>
  <p>{_esc(o.summary)}</p>
  {exists}
  <div class="cols">
    <div><strong>Matched capabilities</strong><ul>{matched}</ul></div>
    <div><strong>Missing capabilities</strong><ul>{missing}</ul></div>
  </div>
  <table class="costs">
    <tr><th>Build (USD)</th><th>Per run</th><th>Monthly ops</th></tr>
    <tr><td>{o.costs.build_cost_usd_low:,.0f}–{o.costs.build_cost_usd_high:,.0f} {_badge('estimate')}</td>
        <td>{o.costs.per_run_cost_usd:,.2f} {_badge('estimate')}</td>
        <td>{o.costs.monthly_operation_usd:,.2f} {_badge('estimate')}</td></tr>
  </table>
  <p class="method"><strong>Method:</strong> {_esc(o.costs.method)}</p>
  <details><summary>Assumptions</summary><ul>{assumptions}</ul></details>
  <p class="scores">{scores}</p>
  <p class="evidence">Evidence: {_esc(", ".join(o.finding_ids) or "none")}</p>
</div>"""


def render_html(case: CaseFile, path: Path | str) -> Path:
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)

    finding_rows = "".join(
        f'<tr data-kind="{_esc(f.kind)}"><td>{_esc(f.id)}</td><td>{_badge(f.kind)}</td>'
        f"<td>{_esc(f.claim)}</td><td>{_esc(f.category)}</td>"
        f'<td><a href="{_esc(f.source.url)}">{_esc(f.source.title or f.source.url)}</a>'
        f'{" ⚑vendor" if f.vendor_claim else ""}{" ✓" if f.source.verified else " (unverified)"}</td>'
        f"<td>{f.confidence:.2f}</td></tr>"
        for f in case.findings
    )

    matrix_rows = []
    for category, options in case.tool_landscape.items():
        for o in options:
            total = sum(o.scores.values())
            matrix_rows.append(
                f"<tr><td>{_esc(o.name)}</td><td>{_esc(category)}</td>"
                f"<td>{o.similarity.index}</td>"
                f"<td>{o.costs.build_cost_usd_low:,.0f}–{o.costs.build_cost_usd_high:,.0f}</td>"
                f"<td>{o.costs.monthly_operation_usd:,.2f}</td><td>{total:g}</td></tr>"
            )

    cards = "".join(
        f"<h3>{_esc(cat)}</h3>" + "".join(_option_card(o) for o in options)
        for cat, options in case.tool_landscape.items()
    )

    suit = ""
    if case.suitability:
        s = case.suitability
        score_rows = "".join(f"<tr><td>{_esc(k)}</td><td>{v:g}/10</td></tr>" for k, v in s.scores.items())
        suit = (f'<section id="suitability"><h2>Suitability verdict</h2>'
                f'<p class="verdict">{_esc(s.verdict)}</p><table>{score_rows}</table>'
                f"<p>{_esc(s.rationale)}</p>"
                f"<p>Cited findings: {_esc(', '.join(s.cited_finding_ids))}</p></section>")

    # the decision LEADS the report — evidence follows for readers who want it
    decision = ""
    if case.decision:
        d = case.decision
        roi = d.roi
        pilot = ""
        if d.pilot_plan:
            p = d.pilot_plan
            pilot = (f"<p><strong>Pilot plan:</strong> {_esc(p.scope)} — {p.duration_weeks} weeks.<br>"
                     f"Success criteria: {_esc('; '.join(p.success_criteria))}<br>"
                     f"Edge cases to test: {_esc('; '.join(p.edge_cases_to_test))}<br>"
                     f"Approvals needed: {_esc('; '.join(p.approvals_needed) or 'none identified')}</p>")
        roi_row = ""
        if roi.monthly_value_usd or roi.hours_saved_per_month:
            roi_row = (f"<table><tr><th>Hours saved / month {_badge('estimate')}</th>"
                       f"<th>Value / month {_badge('estimate')}</th><th>Implementation {_badge('estimate')}</th>"
                       f"<th>Payback {_badge('estimate')}</th></tr>"
                       f"<tr><td>{roi.hours_saved_per_month:g}</td>"
                       f"<td>${roi.monthly_value_usd:,.0f}</td>"
                       f"<td>${roi.implementation_cost_usd_low:,.0f}–${roi.implementation_cost_usd_high:,.0f}</td>"
                       f"<td>{roi.payback_months_low:g}–{roi.payback_months_high:g} months</td></tr></table>"
                       f"<details><summary>ROI assumptions</summary><ul>"
                       + "".join(f"<li>{_esc(a)}</li>" for a in roi.assumptions) + "</ul></details>")
        decision = (f'<section id="decision" style="border:2px solid #1a7f37;border-radius:10px;'
                    f'padding:4px 18px;background:#f0fff4"><h2>Recommendation</h2>'
                    f'<p class="verdict">{_esc(d.action.upper())}'
                    + (f" — {_esc(d.recommended_option)}" if d.recommended_option else "")
                    + f"</p><p>{_esc(d.rationale)}</p>{roi_row}{pilot}"
                    f"<p class=\"evidence\">Evidence: {_esc(', '.join(d.cited_finding_ids))}</p></section>")

    open_qs = "".join(f"<li>{_esc(q)}</li>" for q in case.open_questions) or "<li>none</li>"

    doc = f"""<!doctype html>
<html><head><meta charset="utf-8"><title>Detailed analysis — run {_esc(case.run_id)}</title>
<style>
 body{{font:15px/1.5 system-ui,sans-serif;margin:0;color:#1f2328}}
 main{{max-width:1100px;margin:0 auto;padding:24px}}
 h1{{font-size:1.5em}} h2{{border-bottom:1px solid #d0d7de;padding-bottom:4px;margin-top:36px}}
 table{{border-collapse:collapse;width:100%;margin:8px 0}}
 th,td{{border:1px solid #d0d7de;padding:6px 8px;text-align:left;vertical-align:top}}
 th{{background:#f6f8fa;cursor:pointer}}
 .badge{{color:#fff;border-radius:10px;padding:1px 8px;font-size:.78em}}
 .label{{background:#0969da;color:#fff;border-radius:4px;padding:1px 6px;font-size:.8em}}
 .card{{border:1px solid #d0d7de;border-radius:8px;padding:12px 16px;margin:12px 0}}
 .card h4 .sim{{float:right;color:#0969da}}
 .cols{{display:flex;gap:24px;flex-wrap:wrap}} .cols>div{{flex:1;min-width:220px}}
 .exists{{background:#fff8c5;padding:6px 10px;border-radius:6px}}
 .verdict{{font-size:1.3em;font-weight:700}}
 .filters button{{margin-right:6px;padding:4px 10px;border:1px solid #d0d7de;border-radius:6px;background:#f6f8fa;cursor:pointer}}
 .filters button.on{{background:#0969da;color:#fff}}
 .legend span{{margin-right:12px}}
 .method,.evidence,.scores{{color:#57606a;font-size:.9em}}
 div.scroll{{overflow-x:auto}}
</style></head><body><main>
<h1>{_esc(case.title or 'Opportunity-to-Solution')} — detailed analysis <small>run {_esc(case.run_id)}</small></h1>
<p class="legend">Every number is labeled: {_badge('fact')} directly cited ·
{_badge('estimate')} derived, method shown · {_badge('assumption')} unverified.</p>
{decision}
<section><h2>Problem</h2><p>{_esc(case.problem_statement)}</p>
<p><strong>Stated:</strong> {_esc(case.stated_vs_real.get('stated', ''))}<br>
<strong>Real:</strong> {_esc(case.stated_vs_real.get('real', ''))}</p></section>

<section><h2>Future workflow (validated)</h2><div class="scroll">
<table><tr><th>ID</th><th>Step</th><th>Actor</th><th>System</th><th>Time</th><th>Pain</th><th>Label / rationale</th></tr>
{_workflow_rows(case.future_workflow)}</table></div></section>

<section><h2>Decision matrix</h2><div class="scroll">
<table id="matrix"><thead><tr><th>Option</th><th>Category</th><th>Similarity</th>
<th>Build USD {_badge('estimate')}</th><th>Monthly USD {_badge('estimate')}</th><th>Total score</th></tr></thead>
<tbody>{''.join(matrix_rows)}</tbody></table></div>
<p class="method">Click a header to sort.</p></section>

<section><h2>Tool landscape</h2>{cards}</section>

<section><h2>Findings ({len(case.findings)})</h2>
<div class="filters">Filter: <button data-k="all" class="on">all</button>
<button data-k="fact">fact</button><button data-k="estimate">estimate</button>
<button data-k="assumption">assumption</button></div><div class="scroll">
<table id="findings"><thead><tr><th>ID</th><th>Kind</th><th>Claim</th><th>Category</th><th>Source</th><th>Conf.</th></tr></thead>
<tbody>{finding_rows}</tbody></table></div></section>

<section><h2>Open questions</h2><ul>{open_qs}</ul></section>
{suit}
<script>
document.querySelectorAll('.filters button').forEach(b=>b.onclick=()=>{{
  document.querySelectorAll('.filters button').forEach(x=>x.classList.remove('on'));
  b.classList.add('on');
  const k=b.dataset.k;
  document.querySelectorAll('#findings tbody tr').forEach(r=>
    r.style.display=(k==='all'||r.dataset.kind===k)?'':'none');
}});
document.querySelectorAll('#matrix th').forEach((th,i)=>th.onclick=()=>{{
  const tb=th.closest('table').tBodies[0];
  const asc=th.dataset.asc!=='1'; th.dataset.asc=asc?'1':'0';
  [...tb.rows].sort((a,b)=>{{
    const x=a.cells[i].innerText, y=b.cells[i].innerText;
    const nx=parseFloat(x.replace(/[^0-9.-]/g,'')), ny=parseFloat(y.replace(/[^0-9.-]/g,''));
    const c=(isNaN(nx)||isNaN(ny))?x.localeCompare(y):nx-ny;
    return asc?c:-c;
  }}).forEach(r=>tb.appendChild(r));
}});
</script>
</main></body></html>"""
    path.write_text(doc, encoding="utf-8")
    return path


def render_one_pager(case: CaseFile, path: Path | str) -> Path:
    """Executive one-pager: problem → recommendation → ROI → risks → next step.
    The artifact a stakeholder forwards; everything else is backup."""
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    d, s = case.decision, case.suitability
    top = sorted((o for opts in case.tool_landscape.values() for o in opts),
                 key=lambda o: -(o.similarity.index or 0))[:3]
    risks = []
    for o in top:
        if o.vendor_only:
            risks.append(f"{o.name}: evidence is vendor-only (unverified independently)")
        if o.community_only:
            risks.append(f"{o.name}: anecdote-only evidence")
    for item in case.data_inventory:
        if item.sensitivity in ("confidential", "regulated"):
            risks.append(f"data '{item.name}' is {item.sensitivity} — needs security/compliance review")
    unresolved = [c.field for c in case.captured if c.status == "missing"]
    if unresolved:
        risks.append(f"never captured: {', '.join(unresolved[:5])}")
    roi_html = ""
    if d and (d.roi.monthly_value_usd or d.roi.hours_saved_per_month):
        roi_html = (f"<div class='kpis'><div><b>{d.roi.hours_saved_per_month:g}</b><span>hours/month saved "
                    f"{_badge('estimate')}</span></div><div><b>${d.roi.monthly_value_usd:,.0f}</b>"
                    f"<span>value/month {_badge('estimate')}</span></div>"
                    f"<div><b>{d.roi.payback_months_low:g}–{d.roi.payback_months_high:g}</b>"
                    f"<span>months payback {_badge('estimate')}</span></div></div>")
    next_step = ("Run the pilot: " + d.pilot_plan.scope if d and d.pilot_plan and d.action in ("pilot", "build")
                 else f"Proceed with {d.recommended_option}" if d and d.recommended_option
                 else d.rationale.split(".")[0] if d else "Complete the research run for a recommendation.")
    doc = f"""<!doctype html><html><head><meta charset="utf-8">
<title>{_esc(case.title or case.run_id)} — decision brief</title>
<style>
 body{{font:15px/1.5 system-ui,sans-serif;color:#1f2328;max-width:760px;margin:24px auto;padding:0 16px}}
 h1{{font-size:1.3em;margin-bottom:2px}} .sub{{color:#57606a;font-size:.85em}}
 .rec{{border:2px solid #1a7f37;background:#f0fff4;border-radius:10px;padding:12px 18px;margin:16px 0}}
 .rec .action{{font-size:1.5em;font-weight:800}}
 .kpis{{display:flex;gap:24px;margin:10px 0}} .kpis div{{text-align:center}}
 .kpis b{{font-size:1.4em;display:block}} .kpis span{{font-size:.75em;color:#57606a}}
 .badge{{color:#fff;border-radius:10px;padding:0 7px;font-size:.7em}}
 ul{{margin:4px 0}} h2{{font-size:1em;text-transform:uppercase;letter-spacing:.5px;color:#57606a;margin:18px 0 4px}}
</style></head><body>
<h1>{_esc(case.title or 'AI/automation opportunity')} </h1>
<p class="sub">Decision brief · run {_esc(case.run_id)} · {_esc(case.updated_at[:10])} ·
facts cited, estimates labeled — full evidence in the detailed report</p>
<h2>Problem</h2><p>{_esc(case.problem_statement)}</p>
{f'<div class="rec"><div class="action">{_esc(d.action.upper())}'
 + (f" — {_esc(d.recommended_option)}" if d.recommended_option else "") + f'</div>'
 + f'<p>{_esc(d.rationale)}</p>{roi_html}</div>' if d else
 f'<div class="rec"><div class="action">{_esc(s.verdict.upper() if s else "IN PROGRESS")}</div></div>'}
<h2>Options considered</h2><ul>
{''.join(f"<li><b>{_esc(o.name)}</b> ({_esc(o.category)}) — similarity {o.similarity.index}/100, "
         f"~${o.costs.monthly_operation_usd:,.0f}/month {_badge('estimate')}"
         f"{' — <b>already exists</b>' if o.similarity.existing_solution else ''}</li>" for o in top)
 or '<li>none surfaced</li>'}</ul>
<h2>Risks & open items</h2><ul>{''.join(f'<li>{_esc(r)}</li>' for r in risks) or '<li>none flagged</li>'}</ul>
<h2>Next step</h2><p>{_esc(next_step)}</p>
</body></html>"""
    path.write_text(doc, encoding="utf-8")
    return path


def render_ppt(case: CaseFile, path: Path | str) -> Path | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        print("[reports] python-pptx not installed — skipping PPT overview "
              "(pip install python-pptx)")
        return None

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    blank = prs.slide_layouts[6]

    def slide(title: str, lines: List[str]) -> None:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.9))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        body = s.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.6))
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines[:14]):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(15)

    slide(f"{case.title or 'Opportunity-to-Solution'} — run {case.run_id}",
          [case.problem_statement or "(no problem statement)",
           "", "fact = cited · estimate = derived (method shown) · assumption = unverified"])

    if case.decision:  # the decision comes FIRST
        d = case.decision
        lines = [f"RECOMMENDATION: {d.action.upper()}"
                 + (f" — {d.recommended_option}" if d.recommended_option else ""),
                 d.rationale]
        if d.roi.monthly_value_usd or d.roi.hours_saved_per_month:
            lines.append(f"ROI (estimate): {d.roi.hours_saved_per_month:g} h/month ≈ "
                         f"${d.roi.monthly_value_usd:,.0f}/month · payback "
                         f"{d.roi.payback_months_low:g}–{d.roi.payback_months_high:g} months")
        if d.pilot_plan:
            lines.append(f"Pilot: {d.pilot_plan.scope} ({d.pilot_plan.duration_weeks} weeks) — "
                         + "; ".join(d.pilot_plan.success_criteria[:3]))
        slide("Recommendation", lines)

    slide("Stated vs real problem",
          [f"Stated: {case.stated_vs_real.get('stated', '')}",
           f"Real: {case.stated_vs_real.get('real', '')}",
           f"Evidence: {case.stated_vs_real.get('evidence', '')}"])

    slide("Future workflow",
          [f"{s.id} {s.name} — [{s.label or 'n/a'}] {s.rationale}" for s in case.future_workflow]
          or ["(not mapped)"])

    for category, options in case.tool_landscape.items():
        lines = []
        for o in options:
            flag = " — ALREADY EXISTS" if o.similarity.existing_solution else ""
            lines.append(f"{o.name}: similarity {o.similarity.index}/100, "
                         f"build ${o.costs.build_cost_usd_low:,.0f}–${o.costs.build_cost_usd_high:,.0f} (estimate), "
                         f"monthly ${o.costs.monthly_operation_usd:,.2f} (estimate){flag}")
        slide(f"Landscape — {category}", lines or ["(no options found)"])

    if case.suitability:
        s = case.suitability
        slide("Suitability verdict",
              [f"Verdict: {s.verdict}",
               "Scores: " + ", ".join(f"{k} {v:g}/10" for k, v in s.scores.items()),
               f"Cited findings: {', '.join(s.cited_finding_ids)}",
               (f"Better path: {s.better_path}" if s.better_path else "")])

    slide("Open questions / next steps", case.open_questions or ["none"])
    prs.save(str(path))
    return path
